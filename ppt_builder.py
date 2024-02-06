from pptx import Presentation
import pandas as pd
from pptx.util import Inches, Pt

from . import table_style


class PPTBuilder:
    """Object to build automatic PowerPoint. Currently only supports slides with a title and a table of arbitrary dimensions. More slide types to be added in the future.
    """

    def __init__(self):
        # Store constants
        self.font_size = Pt(15)
        self.template_path = None

        # Hardcode dimensions for now
        self.big_table_dimensions = (
            Inches(0.2), Inches(1.5), Inches(9.6), Inches(5))

        # Create presentation object using the default template.
        self.prs = self.create_new_presentation(None)
        self.prs.save('temp.pptx')
        self.template_path = 'temp.pptx'

        self.styles = {}
        self.content_added = False

    def load_template(self, path: str):
        """Loads PowerPoint template. This can be any PowerPoint (.pptx file)

        Args:
            path (str): Path to template
        """
        self.template_path = path
        self.prs = self.create_new_presentation(path)

    def reset_content(self):
        """Fully reset presentation object. This is required after editing the template. The current template is used, so first update the template
        """
        self.prs = self.create_new_presentation(self.template_path)

    def create_new_presentation(self, template_path: str = None):
        """Generate new Presentation object for the builder

        Args:
            template_path (str, optional): Path to load template from. If None, an empty presentation is initialized.

        Returns:
            Presentation: Presentation object
        """
        if template_path:
            return Presentation(template_path)
        else:
            return Presentation()

    def add_slide(self, template_id: int = 0, title_content: str = ""):
        """Add slide to presentation. Use a specific template if a template id is provided. Template defaults to a title slide.

        Args:
            template_id (int): Template ID. Corresponds to the different Slide Template options in PPT. 
            Different than GUID, this indicates a template slide from the template file.


            title_content (str, optional): Slide Title. Defaults to "".

        Returns:
            pptx Slide: Created slide object (is already added to Presentation)
        """

        slide_layout = self.prs.slide_layouts[template_id]
        slide = self.prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_content

        return slide

    def add_table_style_to_template(self, style: tuple, name: str, include_header_row: bool = True):
        """Add custom Table Style to current template. Colors should be provided in HEX format as strings.
           See https://color-hex.org/ for color reference. The hex code for white would be "FFFFFF".

        Args:
            style (tuple): Style tuple in colors, three colors required. (header color, even row color, odd row color)
            name (name): Style Name
            include_header_row (bool, optional): Whether to include header row. Defaults to True.

        Raises:
            ValueError: If content was modified, raise Error. First add styles, then add content. 
        """

        # Add error if styles are added after content. This is not allowed
        if self.content_added:
            raise ValueError(
                "Do not add custom styles to builder after adding content.")

        # Add styles to template powerpoint
        style_id = table_style.add_style_to_powerpoint(self.template_path, self.template_path, name,
                                                       style[0], style[1], style[2], include_first_row=include_header_row)

        self.styles[name] = style_id

        # Reload template
        self.reset_content()

    def add_table_only_slide(self, df: pd.DataFrame, style_name: str = None, template_id: int = 0, title: str = ""):
        """Add slide to presentation that contains one table and a title

        Args:
            df (pandas DataFrame): Table to be inserted
            style_name (Name, optional): TableStyle name. Defaults to None (no styling).
            template_id (int, optional): Template ID for slide format. Defaults to 0.
            title (str, optional): Slide title. Defaults to "".

        """
        self.content_added = True

        # Create slide
        slide = self.add_slide(template_id=template_id, title_content=title)

        # Get dimensions for table
        num_rows, num_columns = df.shape
        num_rows += 1
        left, top, width, height = self.big_table_dimensions

        # Add the table
        table = slide.shapes.add_table(
            num_rows, num_columns, left, top, width, height)

        # Select the style id to be inserted in table object
        style = self.styles[style_name]

        # Insert style id in table object
        table._element.graphic.graphicData.tbl[0][-1].text = style

        table = table.table

        for i in range(num_columns):
            table.cell(0, i).text = df.columns[i]

        # Add rest of content.
        for index, row in list(df.iterrows()):
            for j in range(num_columns):
                table.cell(index + 1, j).text = row[df.columns[j]]

        # Set font size for all content
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = self.font_size

        # You can adjust the column widths here
        # table.columns[i].width = Inches(width) -> Use to set width of column i
        # 0 is first column, etc.

    def add_title_only_slide(self, template_id: int = 0, title: str = ""):
        """Add slide to presentation that contains only a title

        Args:
            template_id (int, optional): Template ID for slide format. Defaults to 0.
            title (str, optional): Slide title. Defaults to "".
        """
        self.content_added = True
        # Create slide
        self.add_slide(template_id=template_id, title_content=title)

    def save(self, filename: str):
        """Save as powerpoint

        Args:
            filename (str): Location to store powerpoint
        """
        self.prs.save(filename)
